import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from PIL import Image

def createPresentation(groq_content, output_filename, images=None):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    sections = re.split(r"(?=\*\*.*?\*\*)", groq_content.strip())
    if images is None:
        images = []

    def set_font(run, size, bold=False, color=None):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = 'Segoe UI'
        if color:
            run.font.color.rgb = RGBColor(*color)

    def add_slide(content, slide_index, image_url=None):
        slide = prs.slides.add_slide(blank_layout)
        # Minimalist stylish background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 247, 255)

        match = re.match(r"\*\*(.*?)\*\*\s*(.*)", content.strip(), re.DOTALL)
        if match:
            title = match.group(1).strip()
            body_text = match.group(2).strip()
        else:
            print(f"Skipping malformed section: {content[:50]}...")
            return

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8), Inches(1))
        title_tf = title_box.text_frame
        p = title_tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, 34, bold=True, color=(30, 30, 60))
        p.alignment = PP_ALIGN.CENTER
        title_tf.word_wrap = True

        # Alternating layout
        image_on_left = (slide_index % 2 == 0)
        text_width = Inches(4.5)
        image_width = Inches(3.8)
        image_height = Inches(3.8)
        top_margin = Inches(1.6)

        if image_on_left:
            image_left = Inches(0.5)
            text_left = Inches(4.5)
        else:
            text_left = Inches(0.5)
            image_left = Inches(5.0)

        # Add image if available
        if image_url:
            try:
                img_response = requests.get(image_url, timeout=10)
                img_response.raise_for_status()
                img_stream = BytesIO(img_response.content)
                slide.shapes.add_picture(img_stream, image_left, top_margin, image_width, image_height)
            except Exception as e:
                print(f"Image load failed for slide {slide_index}: {e}")

        # Add text
        body_box = slide.shapes.add_textbox(text_left, top_margin, text_width, Inches(4.2))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        paragraphs = [para.strip() for para in body_text.split("•") if para.strip()]
        for para_text in paragraphs:
            p = body_tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {para_text}"
            set_font(run, 22, color=(60, 60, 80))
            p.alignment = PP_ALIGN.LEFT

        # Add subtle line separator for style
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.LINE
            Inches(0.7), Inches(1.4), Inches(8), Pt(1.5)
        )
        line.line.color.rgb = RGBColor(200, 210, 230)

    for i, section in enumerate(sections):
        image_url = images[i] if i < len(images) else None
        add_slide(section, i, image_url=image_url)

    prs.save(output_filename)
